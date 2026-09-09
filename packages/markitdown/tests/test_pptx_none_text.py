#!/usr/bin/env python3 -m pytest
"""Regression test for #1808: PptxConverter must not crash when ``shape.text``
or ``notes_text_frame.text`` returns ``None`` (e.g. an ``<a:r>`` run with no
``<a:t>`` child, or certain third-party-generated decks)."""

import io
import os

import pptx
import pptx.shapes.autoshape
import pptx.slide

from markitdown import MarkItDown

TEST_FILES_DIR = os.path.join(os.path.dirname(__file__), "test_files")
PPTX_FIXTURE = os.path.join(TEST_FILES_DIR, "test.pptx")


def _build_pptx_with_notes() -> io.BytesIO:
    """Build a minimal pptx with a notes slide attached, in memory."""
    prs = pptx.Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.notes_slide.notes_text_frame.text = "some notes"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def test_pptx_converter_handles_none_shape_text(monkeypatch):
    """``shape.text`` returning None must not blow up the title / body paths."""
    monkeypatch.setattr(
        pptx.shapes.autoshape.Shape,
        "text",
        property(lambda self: None),
        raising=True,
    )

    result = MarkItDown().convert(PPTX_FIXTURE)
    assert result.markdown is not None


def test_pptx_converter_handles_none_notes_text(monkeypatch):
    """``notes_text_frame.text`` returning None must not blow up the notes path."""
    real_notes = pptx.slide.NotesSlide.notes_text_frame

    class _NoneText:
        text = None

    def fake_notes(self):
        frame = real_notes.fget(self)
        if frame is None:
            return None
        return _NoneText()

    monkeypatch.setattr(
        pptx.slide.NotesSlide,
        "notes_text_frame",
        property(fake_notes),
        raising=True,
    )

    buf = _build_pptx_with_notes()
    result = MarkItDown().convert_stream(buf, file_extension=".pptx")
    assert result.markdown is not None
