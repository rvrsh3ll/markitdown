#!/usr/bin/env python3 -m pytest
import io
import json
import ntpath
import os
import re
import shutil
import zipfile
import pytest
from types import SimpleNamespace
from typing import Optional
from unittest.mock import MagicMock

import markitdown._uri_utils as uri_utils
from markitdown._uri_utils import parse_data_uri, file_uri_to_path
from markitdown.converters._wikipedia_converter import WikipediaConverter
from markitdown._markitdown import _get_content_disposition_filename
from markitdown.converters import RssConverter

from markitdown import (
    MarkItDown,
    UnsupportedFormatException,
    FileConversionException,
    StreamInfo,
)
from markitdown.converters import YouTubeConverter

# This file contains module tests that are not directly tested by the FileTestVectors.
# This includes things like helper functions and runtime conversion options
# (e.g., LLM clients, exiftool path, transcription services, etc.)

skip_remote = (
    True if os.environ.get("GITHUB_ACTIONS") else False
)  # Don't run these tests in CI


# Don't run the llm tests without a key and the client library
skip_llm = False if os.environ.get("OPENAI_API_KEY") else True
try:
    import openai
except ModuleNotFoundError:
    skip_llm = True

# Skip exiftool tests if not installed
skip_exiftool = shutil.which("exiftool") is None

TEST_FILES_DIR = os.path.join(os.path.dirname(__file__), "test_files")

JPG_TEST_EXIFTOOL = {
    "Author": "AutoGen Authors",
    "Title": "AutoGen: Enabling Next-Gen LLM Applications via Multi-Agent Conversation",
    "Description": "AutoGen enables diverse LLM-based applications",
    "ImageSize": "1615x1967",
    "DateTimeOriginal": "2024:03:14 22:10:00",
}

MP3_TEST_EXIFTOOL = {
    "Title": "f67a499e-a7d0-4ca3-a49b-358bd934ae3e",
    "Artist": "Artist Name Test String",
    "Album": "Album Name Test String",
    "SampleRate": "48000",
}

PDF_TEST_URL = "https://arxiv.org/pdf/2308.08155v2.pdf"
PDF_TEST_STRINGS = [
    "While there is contemporaneous exploration of multi-agent approaches"
]

YOUTUBE_TEST_URL = "https://www.youtube.com/watch?v=V2qZ_lgxTzg"
YOUTUBE_TEST_STRINGS = [
    "## AutoGen FULL Tutorial with Python (Step-By-Step)",
    "This is an intermediate tutorial for installing and using AutoGen locally",
    "PT15M4S",
    "the model we're going to be using today is GPT 3.5 turbo",  # From the transcript
]

DOCX_COMMENT_TEST_STRINGS = [
    "314b0a30-5b04-470b-b9f7-eed2c2bec74a",
    "49e168b7-d2ae-407f-a055-2167576f39a1",
    "## d666f1f7-46cb-42bd-9a39-9a39cf2a509f",
    "# Abstract",
    "# Introduction",
    "AutoGen: Enabling Next-Gen LLM Applications via Multi-Agent Conversation",
    "This is a test comment. 12df-321a",
    "Yet another comment in the doc. 55yiyi-asd09",
]

BLOG_TEST_URL = "https://microsoft.github.io/autogen/blog/2023/04/21/LLM-tuning-math"
BLOG_TEST_STRINGS = [
    "Large language models (LLMs) are powerful tools that can generate natural language texts for various applications, such as chatbots, summarization, translation, and more. GPT-4 is currently the state of the art LLM in the world. Is model selection irrelevant? What about inference parameters?",
    "an example where high cost can easily prevent a generic complex",
]

LLM_TEST_STRINGS = [
    "5bda1dd6",
]

PPTX_TEST_STRINGS = [
    "2cdda5c8-e50e-4db4-b5f0-9722a649f455",
    "04191ea8-5c73-4215-a1d3-1cfb43aaaf12",
    "44bf7d06-5e7a-4a40-a2e1-a2e42ef28c8a",
    "1b92870d-e3b5-4e65-8153-919f4ff45592",
    "AutoGen: Enabling Next-Gen LLM Applications via Multi-Agent Conversation",
    "a3f6004b-6f4f-4ea8-bee3-3741f4dc385f",  # chart title
    "2003",  # chart value
]


def test_wikipedia_converter_no_title() -> None:
    """WikipediaConverter should not render '# None' when page has no title."""
    converter = WikipediaConverter()
    html = b"<html><body><div id='mw-content-text'><p>Hello</p></div></body></html>"
    stream_info = StreamInfo(
        mimetype="text/html", url="https://en.wikipedia.org/wiki/Test"
    )
    result = converter.convert(io.BytesIO(html), stream_info)
    assert "# None" not in result.markdown
    assert "Hello" in result.markdown
    assert result.markdown.strip() == "Hello"


def test_wikipedia_converter_blank_title() -> None:
    """WikipediaConverter should not render an empty heading for a blank title."""
    converter = WikipediaConverter()
    html = b"<html><head><title>   </title></head><body><div id='mw-content-text'><p>Hello</p></div></body></html>"
    stream_info = StreamInfo(
        mimetype="text/html", url="https://en.wikipedia.org/wiki/Test"
    )
    result = converter.convert(io.BytesIO(html), stream_info)
    assert not result.markdown.lstrip().startswith("#")
    assert result.title is None
    assert result.text_content.strip() == "Hello"


# --- Helper Functions ---
def validate_strings(result, expected_strings, exclude_strings=None):
    """Validate presence or absence of specific strings."""
    text_content = result.text_content.replace("\\", "")
    for string in expected_strings:
        assert string in text_content
    if exclude_strings:
        for string in exclude_strings:
            assert string not in text_content


def test_stream_info_operations() -> None:
    """Test operations performed on StreamInfo objects."""

    stream_info_original = StreamInfo(
        mimetype="mimetype.1",
        extension="extension.1",
        charset="charset.1",
        filename="filename.1",
        local_path="local_path.1",
        url="url.1",
    )

    # Check updating all attributes by keyword
    keywords = ["mimetype", "extension", "charset", "filename", "local_path", "url"]
    for keyword in keywords:
        updated_stream_info = stream_info_original.copy_and_update(
            **{keyword: f"{keyword}.2"}
        )

        # Make sure the targeted attribute is updated
        assert getattr(updated_stream_info, keyword) == f"{keyword}.2"

        # Make sure the other attributes are unchanged
        for k in keywords:
            if k != keyword:
                assert getattr(stream_info_original, k) == getattr(
                    updated_stream_info, k
                )

    # Check updating all attributes by passing a new StreamInfo object
    keywords = ["mimetype", "extension", "charset", "filename", "local_path", "url"]
    for keyword in keywords:
        updated_stream_info = stream_info_original.copy_and_update(
            StreamInfo(**{keyword: f"{keyword}.2"})
        )

        # Make sure the targeted attribute is updated
        assert getattr(updated_stream_info, keyword) == f"{keyword}.2"

        # Make sure the other attributes are unchanged
        for k in keywords:
            if k != keyword:
                assert getattr(stream_info_original, k) == getattr(
                    updated_stream_info, k
                )

    # Check mixing and matching
    updated_stream_info = stream_info_original.copy_and_update(
        StreamInfo(extension="extension.2", filename="filename.2"),
        mimetype="mimetype.3",
        charset="charset.3",
    )
    assert updated_stream_info.extension == "extension.2"
    assert updated_stream_info.filename == "filename.2"
    assert updated_stream_info.mimetype == "mimetype.3"
    assert updated_stream_info.charset == "charset.3"
    assert updated_stream_info.local_path == "local_path.1"
    assert updated_stream_info.url == "url.1"

    # Check multiple StreamInfo objects
    updated_stream_info = stream_info_original.copy_and_update(
        StreamInfo(extension="extension.4", filename="filename.5"),
        StreamInfo(mimetype="mimetype.6", charset="charset.7"),
    )
    assert updated_stream_info.extension == "extension.4"
    assert updated_stream_info.filename == "filename.5"
    assert updated_stream_info.mimetype == "mimetype.6"
    assert updated_stream_info.charset == "charset.7"
    assert updated_stream_info.local_path == "local_path.1"
    assert updated_stream_info.url == "url.1"


@pytest.mark.parametrize(
    ("url", "video_id"),
    [
        ("https://www.youtube.com/watch?v=V2qZ_lgxTzg", "V2qZ_lgxTzg"),
        ("https://youtu.be/V2qZ_lgxTzg", "V2qZ_lgxTzg"),
        ("https://www.youtube.com/shorts/V2qZ_lgxTzg", "V2qZ_lgxTzg"),
        ("https://www.youtube.com/embed/V2qZ_lgxTzg", "V2qZ_lgxTzg"),
    ],
)
def test_youtube_converter_extracts_supported_video_ids(
    url: str, video_id: str
) -> None:
    converter = YouTubeConverter()
    assert converter._get_video_id(url) == video_id


@pytest.mark.parametrize(
    "url",
    [
        "https://www.youtube.com/watch?v=V2qZ_lgxTzg",
        "https://youtu.be/V2qZ_lgxTzg",
        "https://www.youtube.com/shorts/V2qZ_lgxTzg",
    ],
)
def test_youtube_converter_accepts_supported_url_formats(url: str) -> None:
    converter = YouTubeConverter()
    assert converter.accepts(
        io.BytesIO(b"<html></html>"),
        StreamInfo(url=url, extension=".html"),
    )


def test_data_uris() -> None:
    # Test basic parsing of data URIs
    data_uri = "data:text/plain;base64,SGVsbG8sIFdvcmxkIQ=="
    mime_type, attributes, data = parse_data_uri(data_uri)
    assert mime_type == "text/plain"
    assert len(attributes) == 0
    assert data == b"Hello, World!"

    data_uri = "data:base64,SGVsbG8sIFdvcmxkIQ=="
    mime_type, attributes, data = parse_data_uri(data_uri)
    assert mime_type is None
    assert len(attributes) == 0
    assert data == b"Hello, World!"

    data_uri = "data:text/plain;charset=utf-8;base64,SGVsbG8sIFdvcmxkIQ=="
    mime_type, attributes, data = parse_data_uri(data_uri)
    assert mime_type == "text/plain"
    assert len(attributes) == 1
    assert attributes["charset"] == "utf-8"
    assert data == b"Hello, World!"

    data_uri = "data:text/plain;CHARSET=utf-8;BASE64,SGVsbG8sIFdvcmxkIQ=="
    mime_type, attributes, data = parse_data_uri(data_uri)
    assert mime_type == "text/plain"
    assert len(attributes) == 1
    assert attributes["charset"] == "utf-8"
    assert data == b"Hello, World!"

    data_uri = "data:,Hello%2C%20World%21"
    mime_type, attributes, data = parse_data_uri(data_uri)
    assert mime_type is None
    assert len(attributes) == 0
    assert data == b"Hello, World!"

    data_uri = "data:text/plain,Hello%2C%20World%21"
    mime_type, attributes, data = parse_data_uri(data_uri)
    assert mime_type == "text/plain"
    assert len(attributes) == 0
    assert data == b"Hello, World!"

    data_uri = "data:text/plain;charset=utf-8,Hello%2C%20World%21"
    mime_type, attributes, data = parse_data_uri(data_uri)
    assert mime_type == "text/plain"
    assert len(attributes) == 1
    assert attributes["charset"] == "utf-8"
    assert data == b"Hello, World!"


def test_uppercase_data_image_uri_is_truncated_by_default() -> None:
    markitdown = MarkItDown()
    html = b'<html><body><img alt="dot" src="DATA:image/png;base64,AAAA"></body></html>'
    stream_info = StreamInfo(mimetype="text/html", extension=".html")

    result = markitdown.convert_stream(io.BytesIO(html), stream_info=stream_info)
    assert result.markdown == "![dot](DATA:image/png;base64...)"
    assert "AAAA" not in result.markdown

    result = markitdown.convert_stream(
        io.BytesIO(html), stream_info=stream_info, keep_data_uris=True
    )
    assert result.markdown == "![dot](DATA:image/png;base64,AAAA)"


def test_file_uris() -> None:
    # Test file URI with an empty host
    file_uri = "file:///path/to/file.txt"
    netloc, path = file_uri_to_path(file_uri)
    assert netloc is None
    assert path == "/path/to/file.txt"

    # Test file URI with no host
    file_uri = "file:/path/to/file.txt"
    netloc, path = file_uri_to_path(file_uri)
    assert netloc is None
    assert path == "/path/to/file.txt"

    # Test file URI with localhost
    file_uri = "file://localhost/path/to/file.txt"
    netloc, path = file_uri_to_path(file_uri)
    assert netloc == "localhost"
    assert path == "/path/to/file.txt"

    # URI schemes are case-insensitive
    file_uri = "FILE:///path/to/file.txt"
    netloc, path = file_uri_to_path(file_uri)
    assert netloc is None
    assert path == "/path/to/file.txt"


def test_convert_case_insensitive_uri_schemes(tmp_path) -> None:
    markitdown = MarkItDown()

    data_result = markitdown.convert("DATA:text/plain;base64,SGVsbG8sIFdvcmxkIQ==")
    assert data_result.markdown == "Hello, World!"

    text_file = tmp_path / "hello.txt"
    text_file.write_text("Hello from file", encoding="utf-8")

    file_result = markitdown.convert(text_file.as_uri().replace("file:", "FILE:", 1))
    assert file_result.markdown == "Hello from file"

    # Test file URI with query parameters
    file_uri = "file:///path/to/file.txt?param=value"
    netloc, path = file_uri_to_path(file_uri)
    assert netloc is None
    assert path == "/path/to/file.txt"

    # Test file URI with fragment
    file_uri = "file:///path/to/file.txt#fragment"
    netloc, path = file_uri_to_path(file_uri)
    assert netloc is None
    assert path == "/path/to/file.txt"


def test_file_uri_with_percent_encoded_windows_drive(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from nturl2path import url2pathname as windows_url2pathname

    monkeypatch.setattr(uri_utils, "os", SimpleNamespace(name="nt", path=ntpath))
    monkeypatch.setattr(uri_utils, "url2pathname", windows_url2pathname)

    netloc, path = uri_utils.file_uri_to_path("file:///C%3A/Temp/example.md")

    assert netloc is None
    assert path == r"C:\Temp\example.md"


def test_docx_comments() -> None:
    # Test DOCX processing, with comments and setting style_map on init
    markitdown_with_style_map = MarkItDown(style_map="comment-reference => ")
    result = markitdown_with_style_map.convert(
        os.path.join(TEST_FILES_DIR, "test_with_comment.docx")
    )
    validate_strings(result, DOCX_COMMENT_TEST_STRINGS)


def _write_underlined_docx(path, embedded_style_map: Optional[str] = None) -> str:
    """Write a minimal .docx holding one underlined run, and return its path."""
    document_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p>
      <w:r><w:t>plain </w:t></w:r>
      <w:r><w:rPr><w:u w:val="single"/></w:rPr><w:t>underlined</w:t></w:r>
    </w:p>
  </w:body>
</w:document>"""

    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "word/_rels/document.xml.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"></Relationships>""",
        )
        archive.writestr("word/document.xml", document_xml)

    if embedded_style_map is not None:
        import mammoth

        with open(path, "r+b") as f:
            mammoth.embed_style_map(f, embedded_style_map)

    return str(path)


def test_docx_underlined_text_is_preserved(tmp_path) -> None:
    docx_file = _write_underlined_docx(tmp_path / "underlined.docx")

    result = MarkItDown().convert(docx_file)

    assert "plain <u>underlined</u>" in result.markdown


def test_docx_embedded_style_map_overrides_underline_default(tmp_path) -> None:
    # A style map embedded in the document takes precedence over the default
    # "u => u" mapping that preserves underlines.
    docx_file = _write_underlined_docx(
        tmp_path / "embedded.docx", embedded_style_map="u => em"
    )

    result = MarkItDown().convert(docx_file)

    assert "plain *underlined*" in result.markdown
    assert "<u>" not in result.markdown


def test_docx_caller_style_map_overrides_embedded_style_map(tmp_path) -> None:
    # ... and a caller-supplied style map still outranks the embedded one.
    docx_file = _write_underlined_docx(
        tmp_path / "embedded.docx", embedded_style_map="u => em"
    )

    result = MarkItDown(style_map="u => strong").convert(docx_file)

    assert "plain **underlined**" in result.markdown


def test_html_strikethrough_variants(tmp_path) -> None:
    html = """<!doctype html>
<html><body>
<p>Plain <s>s element</s> after.</p>
<p>Plain <del>del element</del> after.</p>
<p>Plain <strike>strike element</strike> after.</p>
<p>Spaces A<strike> B </strike>C.</p>
<p>Runs D<strike>  E  </strike>F.</p>
<p>Empty G<strike></strike>H.</p>
<p>Newline I<strike>J
K</strike>L.</p>
<p>Break M<strike>N<br>O</strike>P.</p>
</body></html>
"""
    path = tmp_path / "strike.html"
    path.write_text(html, encoding="utf-8")
    markdown = MarkItDown().convert(str(path)).markdown

    assert markdown == "\n\n".join(
        [
            # <s>, <del> and the obsolete <strike> all mean strikethrough
            "Plain ~~s element~~ after.",
            "Plain ~~del element~~ after.",
            "Plain ~~strike element~~ after.",
            # Surrounding whitespace stays outside of the markup ...
            "Spaces A ~~B~~ C.",
            # ... and runs of it collapse to a single space
            "Runs D ~~E~~ F.",
            # An empty element contributes nothing
            "Empty GH.",
            # A line break inside the element is kept, and the markup
            # survives it because strikethrough may span a single newline
            "Newline I~~J\nK~~L.",
            "Break M~~N\nO~~P.",
        ]
    )


def test_docx_equations() -> None:
    markitdown = MarkItDown()
    docx_file = os.path.join(TEST_FILES_DIR, "equations.docx")
    result = markitdown.convert(docx_file)

    # Check for inline equation m=1 (wrapped with single $) is present
    assert "$m=1$" in result.text_content, "Inline equation $m=1$ not found"

    # Find block equations wrapped with double $$ and check if they are present
    block_equations = re.findall(r"\$\$(.+?)\$\$", result.text_content)
    assert block_equations, "No block equations found in the document."


def test_docx_zip_filename_casing_mismatch() -> None:
    """Test that DOCX files with inconsistent ZIP filename casing are handled.

    Some document generators produce .docx files where the central directory
    records one casing (e.g. 'word/document.xml') but the local file headers
    record another (e.g. 'Word/Document.XML'). Python's zipfile module raises
    BadZipFile when reading such files. This test verifies that MarkItDown
    handles this gracefully.

    See: https://github.com/microsoft/markitdown/issues/1812
    """
    import struct

    markitdown = MarkItDown()
    docx_file = os.path.join(TEST_FILES_DIR, "test.docx")

    # Read the original docx and get its expected content
    original_result = markitdown.convert(docx_file)
    assert original_result.markdown.strip(), "Original DOCX should have content"

    # Read raw bytes and corrupt the local file header filenames
    with open(docx_file, "rb") as f:
        raw = bytearray(f.read())

    # Find all local file headers and uppercase their filenames
    corrupted = bytearray(raw)
    offset = 0
    patched_count = 0
    while offset + 30 <= len(corrupted):
        if corrupted[offset : offset + 4] != b"PK\x03\x04":
            break
        fname_len = struct.unpack_from("<H", corrupted, offset + 26)[0]
        extra_len = struct.unpack_from("<H", corrupted, offset + 28)[0]
        if offset + 30 + fname_len > len(corrupted):
            break
        # Uppercase the filename in the local header
        old_name = corrupted[offset + 30 : offset + 30 + fname_len]
        new_name = old_name.upper()
        if old_name != new_name:
            corrupted[offset + 30 : offset + 30 + fname_len] = new_name
            patched_count += 1
        comp_size = struct.unpack_from("<I", corrupted, offset + 18)[0]
        offset = offset + 30 + fname_len + extra_len + comp_size

    assert patched_count > 0, "Should have patched at least one local file header"

    # Verify the corrupted file would fail with plain zipfile
    with pytest.raises(zipfile.BadZipFile):
        with zipfile.ZipFile(io.BytesIO(bytes(corrupted)), "r") as zf:
            for name in zf.namelist():
                zf.read(name)

    # Verify MarkItDown can still convert it
    corrupted_result = markitdown.convert_stream(
        io.BytesIO(bytes(corrupted)),
        file_extension=".docx",
    )
    assert (
        corrupted_result.markdown.strip()
    ), "Corrupted DOCX should still produce content"
    # Content should be equivalent to the original
    assert (
        original_result.markdown.strip() == corrupted_result.markdown.strip()
    ), "Corrupted DOCX should produce the same output as original"


def test_docx_zip_filename_non_casing_mismatch_still_rejected() -> None:
    """Only case-only local/central filename disagreements are repaired.

    Equal encoded length does not imply two names differ only in case, so the
    repair must not be used to wave through archives whose local and central
    directory names genuinely disagree -- zipfile rejects those for good
    reason, and honouring that keeps the workaround scoped to issue #1812.
    """
    import struct

    from markitdown.converter_utils.docx.pre_process import (
        _fix_zip_filename_casing,
    )

    def build_zip(first_name: str) -> bytes:
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.writestr(first_name, b"payload-a")
            zf.writestr("word/document.xml", b"payload-b")
        return buf.getvalue()

    def replace_first_local_name(data: bytes, new_name: bytes) -> bytes:
        """Rewrite only the first local file header's name, leaving the central
        directory -- the authoritative copy -- untouched."""
        raw = bytearray(data)
        fname_len = struct.unpack_from("<H", raw, 26)[0]
        assert len(new_name) == fname_len, "mutation must preserve the length"
        raw[30 : 30 + fname_len] = new_name
        return bytes(raw)

    def reads_cleanly(data: bytes) -> bool:
        try:
            with zipfile.ZipFile(io.BytesIO(data), "r") as zf:
                for name in zf.namelist():
                    zf.read(name)
        except zipfile.BadZipFile:
            return False
        return True

    original = "[Content_Types].xml"

    # A case-only mismatch is repaired ...
    case_only = replace_first_local_name(build_zip(original), original.upper().encode())
    assert not reads_cleanly(case_only)
    repaired = _fix_zip_filename_casing(io.BytesIO(case_only))
    assert reads_cleanly(repaired.read())

    # ... but an equal-length mismatch that is not case-only is left alone,
    # so zipfile still refuses the archive.
    for mutated_name in (
        b"ZBnoudou^Uxqdr\\/ylm",  # unrelated, same length
        b"[Content_Types].xmy",  # a single differing character
        b"[content_types]/xml",  # differs only outside the cased characters
    ):
        data = replace_first_local_name(build_zip(original), mutated_name)
        assert not reads_cleanly(data), f"{mutated_name!r} should not be readable"
        result = _fix_zip_filename_casing(io.BytesIO(data))
        assert not reads_cleanly(
            result.read()
        ), f"{mutated_name!r} must not be silently repaired"


def test_docx_malformed_equations() -> None:
    """Malformed equations should not crash the converter (issue #1979)."""
    import zipfile
    from io import BytesIO

    docx_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math"
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"
            mc:Ignorable="w14 wp14">
  <w:body>
    <w:p>
      <w:r><w:t>Normal text</w:t></w:r>
    </w:p>
    <m:oMathPara>
      <m:oMath>
        <m:r><m:t>x+1</m:t></m:r>
      </m:oMath>
    </m:oMathPara>
    <w:p>
      <w:r><w:t>After good equation</w:t></w:r>
    </w:p>
    <m:oMathPara>
      <!-- oMathPara with no oMath child -->
    </m:oMathPara>
    <w:p>
      <w:r><w:t>After empty oMathPara</w:t></w:r>
    </w:p>
    <m:oMath>
      <!-- empty inline oMath -->
    </m:oMath>
    <w:p>
      <w:r><w:t>After empty inline oMath</w:t></w:r>
    </w:p>
    <oMathPara>
      <!-- oMath outside the math namespace: BeautifulSoup matches it by local
           name, but the namespaced lookup in _convert_omath_to_latex does not
           find it, so the conversion has nothing to work with -->
      <oMath><r><t>y+2</t></r></oMath>
    </oMathPara>
    <w:p>
      <w:r><w:t>After unnamespaced oMathPara</w:t></w:r>
    </w:p>
    <oMath><r><t>z+3</t></r></oMath>
    <w:p>
      <w:r><w:t>After unnamespaced oMath</w:t></w:r>
    </w:p>
  </w:body>
</w:document>"""

    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("word/document.xml", docx_xml.encode("utf-8"))
        z.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        z.writestr(
            "word/_rels/document.xml.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
</Relationships>""",
        )
        z.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )

    buf.seek(0)
    markitdown = MarkItDown()
    result = markitdown.convert(buf)
    assert "Normal text" in result.markdown
    # A malformed equation must not abort the math pre-processing step, which
    # would silently drop the LaTeX conversion of every other equation too
    assert "$x+1$" in result.markdown or "$$x+1$$" in result.markdown
    assert "After empty oMathPara" in result.markdown
    assert "After empty inline oMath" in result.markdown
    assert "After unnamespaced oMathPara" in result.markdown
    assert "After unnamespaced oMath" in result.markdown


def test_xlsx_legacy_show_zeroes_sheetview(tmp_path) -> None:
    from openpyxl import Workbook

    base_path = tmp_path / "base.xlsx"
    xlsx_path = tmp_path / "legacy_show_zeroes.xlsx"

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet["A1"] = "hello"
    sheet["B1"] = "world"
    # A cell whose text is byte-identical to the malformed attribute. openpyxl stores it
    # as an inline string, so it lands in the very worksheet part that gets repaired.
    sheet["A2"] = ' showZeroes="0" '
    workbook.save(base_path)

    with zipfile.ZipFile(base_path) as source:
        with zipfile.ZipFile(xlsx_path, "w", zipfile.ZIP_DEFLATED) as target:
            for item in source.infolist():
                data = source.read(item.filename)
                if item.filename == "xl/worksheets/sheet1.xml":
                    data = data.replace(
                        b"<sheetView ", b'<sheetView showZeroes="0" ', 1
                    )
                    # Guard the fixture: the sheet view attribute and the cell text must
                    # both live here, or this test stops exercising the repair.
                    assert data.count(b'showZeroes="0"') == 2
                target.writestr(item, data)

    result = MarkItDown().convert(str(xlsx_path))

    assert "## Data" in result.markdown
    assert "hello" in result.markdown
    assert "world" in result.markdown

    # The repair must not reach into the worksheet's data: the cell keeps its text ...
    assert 'showZeroes="0"' in result.markdown
    # ... and no part of the document was silently renamed on the way through.
    assert "showZeros" not in result.markdown


def test_xlsx_show_zeroes_rename_is_scoped_to_sheet_view_tags() -> None:
    from markitdown.converters._xlsx_converter import _rename_show_zeroes_attribute

    worksheet_xml = (
        b"<worksheet><sheetViews>"
        b'<sheetView showZeroes="0" workbookViewId="0"/>'
        b'<sheetView\n\tshowZeroes="1"\ttabSelected="1"/>'
        b"</sheetViews>"
        # openpyxl ignores custom sheet views entirely, so they need no repair
        b'<customSheetViews><customSheetView showZeroes="0"/></customSheetViews>'
        b'<sheetData><row r="1">'
        b'<c r="A1" t="inlineStr"><is><t xml:space="preserve"> showZeroes="0" '
        b"</t></is></c>"
        b'<c r="B1"><f>IF(A1=" showZeroes=","x","y")</f><v>y</v></c>'
        b"</row></sheetData></worksheet>"
    )

    repaired = _rename_show_zeroes_attribute(worksheet_xml)

    # Every <sheetView> start tag is repaired, whatever separates its attributes ...
    assert repaired.count(b"showZeros=") == 2
    assert b'<sheetView showZeros="0" workbookViewId="0"/>' in repaired
    assert b'<sheetView\n\tshowZeros="1"\ttabSelected="1"/>' in repaired

    # ... and nothing outside those start tags is rewritten.
    assert b'<t xml:space="preserve"> showZeroes="0" </t>' in repaired
    assert b'<f>IF(A1=" showZeroes=","x","y")</f>' in repaired
    assert b'<customSheetView showZeroes="0"/>' in repaired


def test_input_as_strings() -> None:
    markitdown = MarkItDown()

    # Test input from a stream
    input_data = b"<html><body><h1>Test</h1></body></html>"
    result = markitdown.convert_stream(io.BytesIO(input_data))
    assert "# Test" in result.text_content

    # Test input with leading blank characters
    input_data = b"   \n\n\n<html><body><h1>Test</h1></body></html>"
    result = markitdown.convert_stream(io.BytesIO(input_data))
    assert "# Test" in result.text_content


def _mock_response(content_disposition: str) -> MagicMock:
    response = MagicMock()
    response.headers = {"content-disposition": content_disposition}
    response.url = "https://example.com/download"
    response.iter_content.return_value = [b"name,value\nalpha,beta\n"]
    response.raise_for_status.return_value = None
    return response


def test_convert_response_uses_rfc5987_content_disposition_filename() -> None:
    markitdown = MarkItDown()
    result = markitdown.convert_response(
        _mock_response("attachment; filename*=UTF-8''data.csv")
    )

    assert result.markdown == "\n".join(
        [
            "| name | value |",
            "| --- | --- |",
            "| alpha | beta |",
        ]
    )


def test_convert_response_prefers_extended_content_disposition_filename() -> None:
    markitdown = MarkItDown()
    result = markitdown.convert_response(
        _mock_response("attachment; filename=fallback.txt; filename*=UTF-8''data.csv")
    )

    assert result.markdown == "\n".join(
        [
            "| name | value |",
            "| --- | --- |",
            "| alpha | beta |",
        ]
    )


def test_get_content_disposition_filename_decodes_rfc5987() -> None:
    assert (
        _get_content_disposition_filename(
            "attachment; filename=fallback.txt; "
            "filename*=UTF-8''d%C3%A1t%C3%A1%2Ecsv"
        )
        == "d\u00e1t\u00e1.csv"
    )


def test_pptx_chart_multi_series_conversion() -> None:
    """Charts with multiple series and many categories must convert correctly.

    Regression test for the slow path in PptxConverter._convert_chart_to_markdown,
    where ``series.values[idx]`` was evaluated inside the (category x series) loop.
    In python-pptx each ``series.values`` access rescans the cached points via
    XPath (O(n) per lookup), so the old code was O(n^2) per series (and rebuilt
    the whole tuple for every category), making large charts extremely slow.

    The values are now materialized once per series. This test builds a chart
    with enough categories that the regressed code path would be pathologically
    slow, and verifies the resulting Markdown table is correct across series.
    """
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    n_categories = 200
    categories = [f"C{i}" for i in range(n_categories)]
    series_a = [float(i) for i in range(n_categories)]
    series_b = [float(i * 2) for i in range(n_categories)]

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Series A", series_a)
    chart_data.add_series("Series B", series_b)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(5),
        chart_data,
    )

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)

    result = MarkItDown().convert_stream(buffer, file_extension=".pptx")
    md = result.markdown

    # Both series headers are present
    assert "Series A" in md
    assert "Series B" in md
    # First and last categories are present (nothing truncated)
    assert "| C0 |" in md
    assert f"| C{n_categories - 1} |" in md
    # A representative row carries the correct value for each series
    assert "| C10 | 10.0 | 20.0 |" in md


def test_deeply_nested_html_fallback() -> None:
    """Large, deeply nested HTML should fall back to plain-text extraction
    instead of silently returning unconverted HTML (issue #1636).

    Note: This test uses sys.setrecursionlimit to guarantee a RecursionError
    regardless of the host environment's default limit, making it deterministic
    across different platforms and CI configurations.
    """
    import sys
    import warnings

    markitdown = MarkItDown()

    # Use a small recursion limit so the test is environment-independent.
    # We restore the original limit in a finally block to avoid side-effects.
    original_limit = sys.getrecursionlimit()
    low_limit = 200  # well below markdownify's traversal depth for depth=500

    # Build HTML with nesting deep enough to trigger RecursionError
    depth = 500
    html = "<html><body>"
    for _ in range(depth):
        html += '<div style="margin-left:10px">'
    html += "<p>Deep content with <b>bold text</b></p>"
    for _ in range(depth):
        html += "</div>"
    html += "</body></html>"

    try:
        sys.setrecursionlimit(low_limit)
        with warnings.catch_warnings(record=True) as w:
            warnings.simplefilter("always")
            result = markitdown.convert_stream(
                io.BytesIO(html.encode("utf-8")),
                file_extension=".html",
            )

            # Should have emitted a warning about the fallback
            recursion_warnings = [x for x in w if "deeply nested" in str(x.message)]
            assert len(recursion_warnings) > 0

    finally:
        sys.setrecursionlimit(original_limit)

    # The output should contain the text content, not raw HTML
    assert "Deep content" in result.markdown
    assert "bold text" in result.markdown
    assert "<div" not in result.markdown
    assert "<p>" not in result.markdown


def test_deeply_nested_rss_item_fallback() -> None:
    """Deeply nested HTML inside an RSS item should fall back to plain-text
    extraction instead of silently embedding raw unconverted HTML in the
    markdown output (same failure class as the HTML converter fix in #1644).

    Note: This test uses sys.setrecursionlimit to guarantee a RecursionError
    regardless of the host environment's default limit, making it deterministic
    across different platforms and CI configurations.
    """
    import sys
    import warnings

    markitdown = MarkItDown()

    # Use a small recursion limit so the test is environment-independent.
    # We restore the original limit in a finally block to avoid side-effects.
    original_limit = sys.getrecursionlimit()
    low_limit = 200  # well below markdownify's traversal depth for depth=500

    # Build an RSS item whose content is deeply nested HTML
    depth = 500
    item_html = ""
    for _ in range(depth):
        item_html += '<div style="margin-left:10px">'
    item_html += "<p>Deep feed content with <b>bold text</b></p>"
    for _ in range(depth):
        item_html += "</div>"

    rss = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<rss version="2.0" '
        'xmlns:content="http://purl.org/rss/1.0/modules/content/">'
        "<channel>"
        "<title>Test Feed</title>"
        "<description>A test feed</description>"
        "<item>"
        "<title>Deep Item</title>"
        f"<content:encoded><![CDATA[{item_html}]]></content:encoded>"
        "</item>"
        "</channel>"
        "</rss>"
    )
    atom = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<feed xmlns="http://www.w3.org/2005/Atom">'
        "<title>Test Feed</title>"
        "<entry>"
        "<title>Deep Entry</title>"
        f'<content type="html"><![CDATA[{item_html}]]></content>'
        "</entry>"
        "</feed>"
    )

    try:
        sys.setrecursionlimit(low_limit)
        with warnings.catch_warnings(record=True) as w:
            warnings.simplefilter("always")
            result = markitdown.convert_stream(
                io.BytesIO(rss.encode("utf-8")),
                file_extension=".rss",
            )

            # Should have emitted a warning about the fallback
            recursion_warnings = [x for x in w if "deeply nested" in str(x.message)]
            assert len(recursion_warnings) > 0

        # strict=True should expose the conversion failure rather than applying
        # the plain-text fallback.
        with pytest.raises(RecursionError):
            RssConverter().convert(
                io.BytesIO(rss.encode("utf-8")),
                StreamInfo(extension=".rss"),
                strict=True,
            )
        with pytest.raises(RecursionError):
            RssConverter().convert(
                io.BytesIO(atom.encode("utf-8")),
                StreamInfo(extension=".atom"),
                strict=True,
            )
    finally:
        sys.setrecursionlimit(original_limit)

    # The output should contain the text content, not raw HTML
    assert "Deep feed content" in result.markdown
    assert "bold text" in result.markdown
    assert "<div" not in result.markdown
    assert "<p>" not in result.markdown


def test_doc_rlink() -> None:
    # Test for: CVE-2025-11849
    markitdown = MarkItDown()

    # Document with rlink
    docx_file = os.path.join(TEST_FILES_DIR, "rlink.docx")

    # Directory containing the target rlink file
    rlink_tmp_dir = os.path.abspath(os.sep + "tmp")

    # Ensure the tmp directory exists
    if not os.path.exists(rlink_tmp_dir):
        pytest.skip(f"Skipping rlink test; {rlink_tmp_dir} directory does not exist.")
        return

    rlink_file_path = os.path.join(rlink_tmp_dir, "test_rlink.txt")
    rlink_content = "de658225-569e-4e3d-9ed2-cfb6abf927fc"
    b64_prefix = (
        "ZGU2NTgyMjUtNTY5ZS00ZTNkLTllZDItY2ZiNmFiZjk"  # base64 prefix of rlink_content
    )

    if os.path.exists(rlink_file_path):
        with open(rlink_file_path, "r", encoding="utf-8") as f:
            existing_content = f.read()
            if existing_content != rlink_content:
                raise ValueError(
                    f"Existing {rlink_file_path} content does not match expected content."
                )
    else:
        with open(rlink_file_path, "w", encoding="utf-8") as f:
            f.write(rlink_content)

    try:
        result = markitdown.convert(docx_file, keep_data_uris=True).text_content
        assert (
            b64_prefix not in result
        )  # Make sure the target file was NOT embedded in the output
    finally:
        os.remove(rlink_file_path)


@pytest.mark.skipif(
    skip_remote,
    reason="do not run tests that query external urls",
)
def test_markitdown_remote() -> None:
    markitdown = MarkItDown()

    # By URL
    result = markitdown.convert(PDF_TEST_URL)
    for test_string in PDF_TEST_STRINGS:
        assert test_string in result.text_content

    # Youtube
    # result = markitdown.convert(YOUTUBE_TEST_URL)
    # for test_string in YOUTUBE_TEST_STRINGS:
    #    assert test_string in result.text_content


@pytest.mark.skipif(
    skip_remote,
    reason="do not run remotely run speech transcription tests",
)
def test_speech_transcription() -> None:
    markitdown = MarkItDown()

    # Test WAV files, MP3 and M4A files
    for file_name in ["test.wav", "test.mp3", "test.m4a"]:
        result = markitdown.convert(os.path.join(TEST_FILES_DIR, file_name))
        result_lower = result.text_content.lower()
        assert (
            ("1" in result_lower or "one" in result_lower)
            and ("2" in result_lower or "two" in result_lower)
            and ("3" in result_lower or "three" in result_lower)
            and ("4" in result_lower or "four" in result_lower)
            and ("5" in result_lower or "five" in result_lower)
        )


def test_exceptions() -> None:
    # Check that an exception is raised when trying to convert an unsupported format
    markitdown = MarkItDown()
    with pytest.raises(UnsupportedFormatException):
        markitdown.convert(os.path.join(TEST_FILES_DIR, "random.bin"))

    # Check that an exception is raised when trying to convert a file that is corrupted
    with pytest.raises(FileConversionException) as exc_info:
        markitdown.convert(
            os.path.join(TEST_FILES_DIR, "random.bin"), file_extension=".pptx"
        )
    assert len(exc_info.value.attempts) == 1
    assert type(exc_info.value.attempts[0].converter).__name__ == "PptxConverter"


def test_pptx_converter_treats_none_llm_caption_as_empty(monkeypatch) -> None:
    from markitdown.converters import _pptx_converter

    calls = 0

    def none_caption(*args, **kwargs):
        nonlocal calls
        calls += 1
        return None

    monkeypatch.setattr(_pptx_converter, "llm_caption", none_caption)

    result = MarkItDown().convert(
        os.path.join(TEST_FILES_DIR, "test.pptx"),
        llm_client=MagicMock(),
        llm_model="test-model",
    )

    assert calls > 0
    assert (
        "![This phrase of the caption is Human-written.](Picture4.jpg)"
        in result.markdown
    )


@pytest.mark.skipif(
    skip_exiftool,
    reason="do not run if exiftool is not installed",
)
def test_markitdown_exiftool() -> None:
    which_exiftool = shutil.which("exiftool")
    assert which_exiftool is not None

    # Test explicitly setting the location of exiftool
    markitdown = MarkItDown(exiftool_path=which_exiftool)
    result = markitdown.convert(os.path.join(TEST_FILES_DIR, "test.jpg"))
    for key in JPG_TEST_EXIFTOOL:
        target = f"{key}: {JPG_TEST_EXIFTOOL[key]}"
        assert target in result.text_content

    # Test setting the exiftool path through an environment variable
    os.environ["EXIFTOOL_PATH"] = which_exiftool
    markitdown = MarkItDown()
    result = markitdown.convert(os.path.join(TEST_FILES_DIR, "test.jpg"))
    for key in JPG_TEST_EXIFTOOL:
        target = f"{key}: {JPG_TEST_EXIFTOOL[key]}"
        assert target in result.text_content

    # Test some other media types
    result = markitdown.convert(os.path.join(TEST_FILES_DIR, "test.mp3"))
    for key in MP3_TEST_EXIFTOOL:
        target = f"{key}: {MP3_TEST_EXIFTOOL[key]}"
        assert target in result.text_content


def test_markitdown_llm_parameters() -> None:
    """Test that LLM parameters are correctly passed to the client."""
    mock_client = MagicMock()
    mock_response = MagicMock()
    mock_response.choices = [
        MagicMock(
            message=MagicMock(
                content="Test caption with red circle and blue square 5bda1dd6"
            )
        )
    ]
    mock_client.chat.completions.create.return_value = mock_response

    test_prompt = "You are a professional test prompt."
    markitdown = MarkItDown(
        llm_client=mock_client, llm_model="gpt-4o", llm_prompt=test_prompt
    )

    # Test image file
    markitdown.convert(os.path.join(TEST_FILES_DIR, "test_llm.jpg"))

    # Verify the prompt was passed to the OpenAI API
    assert mock_client.chat.completions.create.called
    call_args = mock_client.chat.completions.create.call_args
    messages = call_args[1]["messages"]
    assert len(messages) == 1
    assert messages[0]["content"][0]["text"] == test_prompt

    # Reset the mock for the next test
    mock_client.chat.completions.create.reset_mock()

    # TODO: may only use one test after the llm caption method duplicate has been removed:
    # https://github.com/microsoft/markitdown/pull/1254
    # Test PPTX file
    markitdown.convert(os.path.join(TEST_FILES_DIR, "test.pptx"))

    # Verify the prompt was passed to the OpenAI API for PPTX images too
    assert mock_client.chat.completions.create.called
    call_args = mock_client.chat.completions.create.call_args
    messages = call_args[1]["messages"]
    assert len(messages) == 1
    assert messages[0]["content"][0]["text"] == test_prompt


@pytest.mark.skipif(
    skip_llm,
    reason="do not run llm tests without a key",
)
def test_markitdown_llm() -> None:
    client = openai.OpenAI()
    markitdown = MarkItDown(llm_client=client, llm_model="gpt-4o")

    result = markitdown.convert(os.path.join(TEST_FILES_DIR, "test_llm.jpg"))
    for test_string in LLM_TEST_STRINGS:
        assert test_string in result.text_content

    # This is not super precise. It would also accept "red square", "blue circle",
    # "the square is not blue", etc. But it's sufficient for this test.
    for test_string in ["red", "circle", "blue", "square"]:
        assert test_string in result.text_content.lower()

    # Images embedded in PPTX files
    result = markitdown.convert(os.path.join(TEST_FILES_DIR, "test.pptx"))
    # LLM Captions are included
    for test_string in LLM_TEST_STRINGS:
        assert test_string in result.text_content
    # Standard alt text is included
    validate_strings(result, PPTX_TEST_STRINGS)


def test_pptx_chart_no_title_text_frame() -> None:
    from markitdown.converters._pptx_converter import PptxConverter

    # python-pptx's ChartTitle.text_frame is destructive -- it creates a text
    # frame if one isn't already present, so it never returns None.
    # has_text_frame is the property that actually reflects presence/absence,
    # which is what the has_title=True-but-no-text-frame case looks like
    # against the real library.
    mock_chart = MagicMock()
    mock_chart.has_title = True
    mock_chart.chart_title.has_text_frame = False

    mock_category = MagicMock()
    mock_category.label = "Cat 1"
    mock_chart.plots = [MagicMock(categories=[mock_category])]

    mock_series = MagicMock()
    mock_series.name = "Series 1"
    mock_series.values = [10.0]
    mock_chart.series = [mock_series]

    converter = PptxConverter()
    result = converter._convert_chart_to_markdown(mock_chart)

    assert "### Chart" in result
    assert "Cat 1" in result
    assert "Series 1" in result
    assert ":" not in result


def test_pptx_chart_with_title_text_frame() -> None:
    from markitdown.converters._pptx_converter import PptxConverter

    mock_chart = MagicMock()
    mock_chart.has_title = True
    mock_chart.chart_title.has_text_frame = True
    mock_chart.chart_title.text_frame.text = "Revenue"

    mock_category = MagicMock()
    mock_category.label = "Cat 1"
    mock_chart.plots = [MagicMock(categories=[mock_category])]

    mock_series = MagicMock()
    mock_series.name = "Series 1"
    mock_series.values = [10.0]
    mock_chart.series = [mock_series]

    converter = PptxConverter()
    result = converter._convert_chart_to_markdown(mock_chart)

    assert "### Chart: Revenue" in result


def test_youtube_converter_missing_title_metadata() -> None:
    """Test that YouTubeConverter converts streams with and without title metadata without raising AssertionError."""
    from unittest.mock import patch
    from markitdown.converters._youtube_converter import YouTubeConverter

    converter = YouTubeConverter()
    stream_info = StreamInfo(
        mimetype="text/html",
        extension=".html",
        url="https://www.youtube.com/watch?v=12345",
    )

    with patch(
        "markitdown.converters._youtube_converter.IS_YOUTUBE_TRANSCRIPT_CAPABLE",
        False,
    ):
        # Case 1: Stream with no title metadata or title tag
        html_content_no_title = b"<html><head></head><body>Video Content</body></html>"
        stream_no_title = io.BytesIO(html_content_no_title)
        result_no_title = converter.convert(stream_no_title, stream_info)
        assert result_no_title.title == ""
        assert "# YouTube" in result_no_title.markdown

        # Case 2: Stream with an empty <title> tag
        html_content_empty_title = (
            b"<html><head><title></title></head><body>Video Content</body></html>"
        )
        stream_empty_title = io.BytesIO(html_content_empty_title)
        result_empty_title = converter.convert(stream_empty_title, stream_info)
        assert result_empty_title.title == ""
        assert "# YouTube" in result_empty_title.markdown

        # Case 3: Stream whose title is only available from the <title> tag
        html_content_title_tag = b"<html><head><title>Fallback Title</title></head><body>Video Content</body></html>"
        stream_title_tag = io.BytesIO(html_content_title_tag)
        result_title_tag = converter.convert(stream_title_tag, stream_info)
        assert result_title_tag.title == "Fallback Title"
        assert "# YouTube" in result_title_tag.markdown


def test_zip_stream_no_filename_header() -> None:
    """Regression test: ZipConverter must not render the literal string 'None'
    in the output header when the stream has no associated URL, local path, or
    filename (e.g. when called via convert_stream() without stream_info)."""
    markitdown = MarkItDown()

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("hello.txt", "Hello world")
    buf.seek(0)

    result = markitdown.convert_stream(
        buf, stream_info=StreamInfo(mimetype="application/zip")
    )
    assert result.markdown.startswith("Content from the zip file `(unknown)`:\n\n")
    assert "Hello world" in result.markdown


def test_ipynb_heading_title_preserves_leading_hash() -> None:
    """Heading marker removal must not eat a leading '#' from the title text.

    Regression for https://github.com/microsoft/markitdown/issues/2367
    """
    from markitdown.converters._ipynb_converter import IpynbConverter

    notebook = {
        "nbformat": 4,
        "nbformat_minor": 5,
        "metadata": {},
        "cells": [
            {
                "cell_type": "markdown",
                "source": ["# #hashtag campaign results\n"],
                "metadata": {},
            }
        ],
    }
    result = IpynbConverter()._convert(notebook)
    assert result.title == "#hashtag campaign results"


def test_ipynb_accepts_non_ascii() -> None:
    """IpynbConverter.accepts() must not raise on non-ASCII binary content."""
    from markitdown.converters._ipynb_converter import IpynbConverter

    converter = IpynbConverter()

    # Binary content that is not valid UTF-8 (simulates non-ASCII file)
    binary_data = b"\x80\x81\x82\x83"
    stream_info = StreamInfo(mimetype="application/json", charset="utf-8")

    # Should return False without raising
    result = converter.accepts(io.BytesIO(binary_data), stream_info)
    assert result is False

    # French PDF content (UTF-8 bytes that would crash if decoded as ASCII)
    french_bytes = "lettre d'information sur l'événement".encode("utf-8")
    stream_info_ascii = StreamInfo(mimetype="application/json", charset="ascii")

    result = converter.accepts(io.BytesIO(french_bytes), stream_info_ascii)
    assert result is False

    # Valid notebook content should still be accepted
    notebook_bytes = b'{"nbformat": 4, "nbformat_minor": 5, "cells": []}'
    stream_info_json = StreamInfo(mimetype="application/json", charset="utf-8")

    result = converter.accepts(io.BytesIO(notebook_bytes), stream_info_json)
    assert result is True


def test_epub_metadata_nodevalue():
    from defusedxml.minidom import parseString
    from markitdown.converters._epub_converter import EpubConverter

    xml_data = (
        '<package xmlns:dc="http://purl.org/dc/elements/1.1/">'
        "<dc:title><span>Structured</span> Title</dc:title>"
        "<dc:creator><name>Author 1</name></dc:creator>"
        "<dc:creator>Author 2</dc:creator>"
        "<dc:publisher></dc:publisher>"
        "<dc:description/>"
        "</package>"
    )
    dom = parseString(xml_data)
    converter = EpubConverter()

    title = converter._get_text_from_node(dom, "dc:title")
    assert title == "Structured Title"

    creators = converter._get_all_texts_from_nodes(dom, "dc:creator")
    assert creators == ["Author 1", "Author 2"]

    publisher = converter._get_text_from_node(dom, "dc:publisher")
    assert publisher is None

    missing = converter._get_text_from_node(dom, "dc:date")
    assert missing is None


def test_json_with_late_non_ascii_character(tmp_path) -> None:
    payload = {
        "record": {
            "title": "Example record",
            "abstract": "This is sample test. " * 500,
        },
        "notes": "non-ASCII character: Ã¨",
    }
    json_path = tmp_path / "input.json"
    json_path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    result = MarkItDown().convert(str(json_path))

    assert "non-ASCII character: Ã¨" in result.text_content


###############################################################################
# CSV converter.
###############################################################################


def _convert_csv(data: bytes, charset: str | None = None) -> str:
    stream_info = StreamInfo(extension=".csv", charset=charset)
    return (
        MarkItDown()
        .convert_stream(io.BytesIO(data), stream_info=stream_info)
        .text_content
    )


def test_csv_utf8_bom_is_stripped_from_the_header() -> None:
    result = _convert_csv(b"\xef\xbb\xbfname,age\nAlice,30\n")

    assert "\ufeff" not in result
    assert result.startswith("| name | age |")
    assert "| Alice | 30 |" in result


def test_csv_utf8_bom_is_stripped_when_charset_is_known() -> None:
    result = _convert_csv(b"\xef\xbb\xbfname,age\nAlice,30\n", charset="utf-8")

    assert "\ufeff" not in result
    assert result.startswith("| name | age |")


def test_csv_leading_blank_line_does_not_destroy_the_table() -> None:
    result = _convert_csv(b"\nname,age\nAlice,30\n")

    assert result == "| name | age |\n| --- | --- |\n| Alice | 30 |"


def test_csv_trailing_blank_lines_are_skipped() -> None:
    result = _convert_csv(b"name,age\nAlice,30\n\n\n")

    assert result == "| name | age |\n| --- | --- |\n| Alice | 30 |"


def test_csv_blank_lines_between_rows_are_kept() -> None:
    result = _convert_csv(b"name,age\n\nAlice,30\n\nBob,40\n")

    assert (
        result == "| name | age |\n| --- | --- |\n| Alice | 30 |\n|  |  |\n| Bob | 40 |"
    )


def test_csv_all_blank_input_returns_empty_markdown() -> None:
    result = _convert_csv(b"\n\n\n")

    assert result == ""


@pytest.mark.parametrize(
    "data,expected",
    [
        (
            b"banner\nname,age,city\nAlice,30,Seattle\n",
            "| banner |  |  |\n| --- | --- | --- |\n"
            "| name | age | city |\n| Alice | 30 | Seattle |",
        ),
        (
            b"name,age\nAlice\n\nBob,40,Seattle\nCarol,25\n",
            "| name | age |  |\n| --- | --- | --- |\n"
            "| Alice |  |  |\n|  |  |  |\n"
            "| Bob | 40 | Seattle |\n| Carol | 25 |  |",
        ),
        (
            b"name\nAlice,,\n",
            "| name |  |  |\n| --- | --- | --- |\n| Alice |  |  |",
        ),
        (
            b"name,age,city\nAlice,30\nBob\n",
            "| name | age | city |\n| --- | --- | --- |\n"
            "| Alice | 30 |  |\n| Bob |  |  |",
        ),
    ],
    ids=["preamble", "widest-row-late", "trailing-empty-fields", "widest-header"],
)
def test_csv_table_matches_widest_row(data: bytes, expected: str) -> None:
    assert _convert_csv(data) == expected


def test_csv_pipe_in_cell_is_escaped() -> None:
    result = _convert_csv(b'name,description\nWidget,"cheap | fast"\n')

    # The pipe must be escaped rather than emitted raw, or it splits the row.
    assert "| Widget | cheap \\| fast |" in result


def test_csv_pipe_preceded_by_backslash_is_still_escaped() -> None:
    # `left\|right` must not become `left\\|right`: the doubled backslash is a
    # literal backslash, which would leave the pipe acting as a delimiter and
    # let a renderer drop `right`. The run of backslashes is doubled instead,
    # so the escaping `\|` survives.
    result = _convert_csv(b'name,description\nWidget,"left\\|right"\n')

    assert r"| Widget | left\\\|right |" in result


def test_csv_pipe_preceded_by_two_backslashes_is_still_escaped() -> None:
    # An even-length run is just as dangerous once the naive `\|` is appended,
    # so check a longer run as well.
    result = _convert_csv(b'name,description\nWidget,"left\\\\|right"\n')

    assert r"| Widget | left\\\\\|right |" in result


def test_csv_newline_in_quoted_cell_does_not_split_the_row() -> None:
    result = _convert_csv(b'name,notes\nWidget,"line one\nline two"\n')

    # The table must stay on one line per record; the embedded break collapses
    # to a space.
    assert len(result.splitlines()) == 3
    assert "| Widget | line one line two |" in result


def test_csv_carriage_returns_in_quoted_cell_collapse_to_spaces() -> None:
    result = _convert_csv(b'name,notes\nWidget,"line one\r\nline two\rline three"\n')

    assert len(result.splitlines()) == 3
    assert "| Widget | line one line two line three |" in result


def test_csv_pipe_in_header_is_escaped() -> None:
    result = _convert_csv(b'"a | b",c\n1,2\n')

    assert "| a \\| b | c |" in result


def test_csv_plain_values_are_unchanged() -> None:
    # Guards against over-escaping ordinary content.
    result = _convert_csv(b"name,description\nWidget,cheap and fast\n")

    assert "| Widget | cheap and fast |" in result
    assert "\\" not in result


def test_csv_backslash_without_a_pipe_is_left_alone() -> None:
    # Only backslashes that guard a pipe are doubled; a Windows path stays
    # readable.
    result = _convert_csv(b"name,path\nWidget,C:\\temp\\file.txt\n")

    assert r"| Widget | C:\temp\file.txt |" in result


# ---------------------------------------------------------------------------
# Regression test for issue #1960:
# exiftool_path pointing to a nonexistent binary used to leak a raw
# FileNotFoundError. It should now be wrapped in RuntimeError with a
# message that includes the path.
# ---------------------------------------------------------------------------


def test_exiftool_metadata_with_nonexistent_binary():
    """#1960: nonexistent exiftool_path raises RuntimeError, not FileNotFoundError."""
    from markitdown.converters._exiftool import exiftool_metadata

    exiftool_path = "/this/does/not/exist/exiftool"
    with pytest.raises(
        RuntimeError,
        match=re.escape(f"Failed to invoke exiftool at {exiftool_path}"),
    ):
        exiftool_metadata(io.BytesIO(b""), exiftool_path=exiftool_path)


def test_exiftool_metadata_with_no_path():
    """Sanity check: exiftool_path=None still returns {} (early return)."""
    from markitdown.converters._exiftool import exiftool_metadata

    assert exiftool_metadata(io.BytesIO(b""), exiftool_path=None) == {}


def test_exiftool_metadata_invocation_oserror():
    """#1960: an OSError while running exiftool is wrapped, and the stream is restored."""
    from unittest.mock import patch

    from markitdown.converters._exiftool import exiftool_metadata

    def fake_run(args, **kwargs):
        # The version check succeeds ...
        if "-ver" in args:
            return SimpleNamespace(stdout="12.24\n")
        # ... but the actual metadata invocation fails.
        raise OSError("exiftool vanished")

    exiftool_path = "/usr/bin/exiftool"
    file_stream = io.BytesIO(b"0123456789")
    file_stream.seek(4)

    with patch("markitdown.converters._exiftool.subprocess.run", side_effect=fake_run):
        with pytest.raises(
            RuntimeError,
            match=re.escape(f"Failed to invoke exiftool at {exiftool_path}"),
        ):
            exiftool_metadata(file_stream, exiftool_path=exiftool_path)

    assert file_stream.tell() == 4


if __name__ == "__main__":
    """Runs this file's tests from the command line."""
    for test in [
        test_stream_info_operations,
        test_data_uris,
        test_file_uris,
        test_docx_comments,
        test_docx_zip_filename_casing_mismatch,
        test_docx_zip_filename_non_casing_mismatch_still_rejected,
        test_input_as_strings,
        test_markitdown_remote,
        test_speech_transcription,
        test_exceptions,
        test_doc_rlink,
        test_markitdown_exiftool,
        test_markitdown_llm_parameters,
        test_markitdown_llm,
        test_pptx_chart_no_title_text_frame,
        test_pptx_chart_with_title_text_frame,
        test_ipynb_accepts_non_ascii,
        test_epub_metadata_nodevalue,
        test_exiftool_metadata_with_nonexistent_binary,
        test_exiftool_metadata_with_no_path,
        test_exiftool_metadata_invocation_oserror,
    ]:
        print(f"Running {test.__name__}...", end="")
        test()
        print("OK")
    print("All tests passed!")
