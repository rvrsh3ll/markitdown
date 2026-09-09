#!/usr/bin/env python3 -m pytest
"""A slide whose notes slide carries no text must not get a "Notes:" heading."""

import io

import pptx

from markitdown import MarkItDown, StreamInfo


def _build_pptx(notes: str | None) -> io.BytesIO:
    """Build a one-slide deck in memory, optionally attaching a notes slide."""
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide title"
    if notes is not None:
        # Touching notes_slide creates the part, which is what PowerPoint does
        # for a deck whose notes pane has been opened.
        slide.notes_slide.notes_text_frame.text = notes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _convert(notes: str | None) -> str:
    return (
        MarkItDown()
        .convert_stream(_build_pptx(notes), stream_info=StreamInfo(extension=".pptx"))
        .markdown
    )


def test_empty_notes_slide_produces_no_notes_heading() -> None:
    markdown = _convert("")

    assert "Slide title" in markdown
    assert "### Notes:" not in markdown


def test_whitespace_only_notes_produce_no_notes_heading() -> None:
    markdown = _convert("   \n\n  ")

    assert "Slide title" in markdown
    assert "### Notes:" not in markdown


def test_slide_without_a_notes_slide_produces_no_notes_heading() -> None:
    markdown = _convert(None)

    assert "Slide title" in markdown
    assert "### Notes:" not in markdown


def test_notes_with_text_are_still_emitted() -> None:
    markdown = _convert("Remember to mention the budget.")

    assert "### Notes:\nRemember to mention the budget." in markdown
